"""Build a styled .pptx briefing deck from the same structured JSON input
used by build_docx_report.py.

Both scripts read one JSON file so the report and the deck stay in sync —
you analyze the scraped articles once, and both outputs are just different
renderings of the same "sections" list. A section that should appear in
only one of the two outputs (e.g. a long "Sources" list that suits the
report but clutters a deck) can be excluded with "include_in".

Usage:
    python build_pptx_deck.py --input briefing.json --output briefing.pptx [--color "#2E86AB"] [--template reference.pptx]

JSON input schema: see build_docx_report.py's module docstring. Relevant
per-section keys for the deck:
  - "heading"   -> slide title
  - "paragraphs" + "bullets" -> slide body bullets (paragraphs become
    bullets too; a deck has no separate prose flow like a document does)
  - "image"     -> picture placed on the right half of the slide
  - "source"    -> small caption under the title
  - "include_in": ["pptx"] / ["docx"] -> restrict which output renders it

--template: path to an existing .pptx to use as the visual template (its
theme, fonts, and slide master/layouts carry over). Any slides already in
that file are stripped before new content is added -- only its layouts are
reused. The deck's own slide dimensions and a blank-ish layout are picked
up from the template automatically; without --template, behavior is
unchanged (blank default presentation, fixed 13.333x7.5in size).
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from image_utils import fit_within_box, prepare_image_for_embed, resolve_image_path

FONT_NAME = "맑은 고딕"
TITLE_SLIDE_TITLE_SIZE = Pt(40)
TITLE_SLIDE_SUBTITLE_SIZE = Pt(20)
SLIDE_TITLE_SIZE = Pt(30)
CAPTION_SIZE = Pt(12)
BODY_SIZE = Pt(18)
DEFAULT_COLOR = "1F4E79"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.5)
IMAGE_MAX_WIDTH_IN = 5.6
IMAGE_MAX_HEIGHT_IN = 5.7
IMAGE_WIDTH = Inches(IMAGE_MAX_WIDTH_IN)

HEX_LENGTH = 6
BLANK_LAYOUT_NAME_HINTS = ("blank", "빈 화면", "빈화면")


def default_dims() -> dict:
    return {
        "slide_width": SLIDE_WIDTH,
        "slide_height": SLIDE_HEIGHT,
        "margin": MARGIN,
        "image_width": IMAGE_WIDTH,
        "image_max_width_in": IMAGE_MAX_WIDTH_IN,
        "image_max_height_in": IMAGE_MAX_HEIGHT_IN,
    }


def dims_for_scale(scale: float) -> dict:
    """Scale every layout constant by `scale` (new_width / default_width)
    so a template with different slide dimensions still gets the same
    proportions instead of overflowing or leaving dead space."""
    return {
        "slide_width": Emu(round(SLIDE_WIDTH * scale)),
        "slide_height": Emu(round(SLIDE_HEIGHT * scale)),
        "margin": Emu(round(MARGIN * scale)),
        "image_width": Emu(round(IMAGE_WIDTH * scale)),
        "image_max_width_in": IMAGE_MAX_WIDTH_IN * scale,
        "image_max_height_in": IMAGE_MAX_HEIGHT_IN * scale,
    }


def find_blank_layout(prs: Presentation):
    """Prefer a layout whose name suggests blank/title-only, since layout
    index 6 may not be blank in an arbitrary template. Falls back to the
    last layout in the master's list."""
    layouts = list(prs.slide_masters[0].slide_layouts)
    for layout in layouts:
        name = (layout.name or "").strip().lower()
        if any(hint in name for hint in BLANK_LAYOUT_NAME_HINTS):
            return layout
    return layouts[-1]


def strip_existing_slides(prs: Presentation) -> None:
    """Remove every slide already present in a template file, keeping only
    its theme/master/layouts. python-pptx has no public API for this, so
    drop each slide's relationship id from the slide list and drop the
    now-unused part, per the standard python-pptx recipe."""
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    for sldId in slide_ids:
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def parse_color(value: str) -> RGBColor:
    hex_value = value.lstrip("#")
    if len(hex_value) != HEX_LENGTH or any(
        c not in "0123456789ABCDEFabcdef" for c in hex_value
    ):
        raise ValueError(
            f"'{value}' is not a valid hex color (expected format: RRGGBB or #RRGGBB)"
        )
    return RGBColor.from_string(hex_value.upper())


def set_font(run, size=None, color: RGBColor = None, bold: bool = None, italic: bool = None):
    run.font.name = FONT_NAME
    if size is not None:
        run.font.size = size
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    return run


def blank_slide(prs: Presentation, layout):
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box.text_frame


def add_title_slide(prs: Presentation, layout, title_text: str, subtitle_text: str, color: RGBColor, dims: dict):
    slide = blank_slide(prs, layout)
    slide_width = dims["slide_width"]
    margin = dims["margin"]

    tf = add_textbox(slide, margin, Inches(2.6), slide_width - 2 * margin, Inches(1.6))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    set_font(run, size=TITLE_SLIDE_TITLE_SIZE, color=color, bold=True)

    if subtitle_text:
        sub_tf = add_textbox(slide, margin, Inches(4.2), slide_width - 2 * margin, Inches(0.8))
        sp = sub_tf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        srun = sp.add_run()
        srun.text = subtitle_text
        set_font(srun, size=TITLE_SLIDE_SUBTITLE_SIZE, color=RGBColor(0x59, 0x59, 0x59), italic=True)

    return slide


def add_content_slide(prs: Presentation, layout, section: dict, color: RGBColor, base_dir: Path, dims: dict):
    slide = blank_slide(prs, layout)
    slide_width = dims["slide_width"]
    slide_height = dims["slide_height"]
    margin = dims["margin"]
    image_width = dims["image_width"]
    has_image = bool(section.get("image"))
    body_width = (slide_width - 3 * margin - image_width) if has_image else (slide_width - 2 * margin)

    # Title
    title_tf = add_textbox(slide, margin, Inches(0.4), slide_width - 2 * margin, Inches(1.0))
    tp = title_tf.paragraphs[0]
    trun = tp.add_run()
    trun.text = section.get("heading", "")
    set_font(trun, size=SLIDE_TITLE_SIZE, color=color, bold=True)

    top = Inches(1.3)

    source = section.get("source")
    if source:
        cap_tf = add_textbox(slide, margin, top, body_width, Inches(0.4))
        cp = cap_tf.paragraphs[0]
        crun = cp.add_run()
        crun.text = source
        set_font(crun, size=CAPTION_SIZE, color=RGBColor(0x80, 0x80, 0x80), italic=True)
        top = top + Inches(0.45)

    body_lines = list(section.get("paragraphs", [])) + list(section.get("bullets", []))
    if body_lines:
        body_tf = add_textbox(slide, margin, top, body_width, slide_height - top - margin)
        for i, line in enumerate(body_lines):
            p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            p.text = f"• {line}"
            for run in p.runs:
                set_font(run, size=BODY_SIZE)
            p.space_after = Pt(10)

    if has_image:
        image_path = section["image"]
        resolved = resolve_image_path(image_path, base_dir)
        image_left = slide_width - margin - image_width
        if resolved.exists():
            embed_path, width_px, height_px = prepare_image_for_embed(resolved, base_dir)
            width_in, height_in = fit_within_box(
                width_px, height_px, dims["image_max_width_in"], dims["image_max_height_in"]
            )
            slide.shapes.add_picture(
                str(embed_path), image_left, Inches(1.3), width=Inches(width_in), height=Inches(height_in)
            )
        else:
            warn_tf = add_textbox(slide, image_left, Inches(1.3), image_width, Inches(0.6))
            wp = warn_tf.paragraphs[0]
            wrun = wp.add_run()
            wrun.text = f"[Screenshot not found: {image_path}]"
            set_font(wrun, size=CAPTION_SIZE, italic=True)

    return slide


def section_applies(section: dict, output: str) -> bool:
    include_in = section.get("include_in")
    return include_in is None or output in include_in


def build_deck(
    data: dict, base_dir: Path, color_override: RGBColor = None, template_path: Path = None
) -> Presentation:
    title_text = data.get("title", "Untitled Briefing")
    subtitle_text = data.get("subtitle", "")
    main_color = color_override or parse_color(data.get("color", DEFAULT_COLOR))

    if template_path is not None:
        prs = Presentation(str(template_path))
        strip_existing_slides(prs)
        scale = prs.slide_width / SLIDE_WIDTH
        dims = dims_for_scale(scale)
        layout = find_blank_layout(prs)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        dims = default_dims()
        layout = prs.slide_layouts[6]  # blank layout

    add_title_slide(prs, layout, title_text, subtitle_text, main_color, dims)

    for section in data.get("sections", []):
        if not section_applies(section, "pptx"):
            continue
        add_content_slide(prs, layout, section, main_color, base_dir, dims)

    return prs


def parse_args():
    parser = argparse.ArgumentParser(description="Build a styled .pptx briefing deck from a JSON summary.")
    parser.add_argument("--input", type=Path, required=True, help="Path to the input JSON file")
    parser.add_argument("--output", type=Path, default=Path("briefing.pptx"), help="Output .pptx path")
    parser.add_argument(
        "--color",
        type=parse_color,
        default=None,
        help="Main color as hex code, e.g. #2E86AB (overrides the JSON's 'color' field)",
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=None,
        help="Path to an existing .pptx to use as the visual template (theme/fonts/layouts). "
        "Its own existing slides are stripped before new content is added.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    data = json.loads(args.input.read_text(encoding="utf-8"))
    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation = build_deck(
        data, base_dir=args.input.parent, color_override=args.color, template_path=args.template
    )
    presentation.save(args.output)
    print(f"Deck saved to {args.output.resolve()}")


if __name__ == "__main__":
    main()
